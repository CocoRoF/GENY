"""Document editing tools — pptx / xlsx / docx edit + convert + preview.

workspace-canvas P3 (docs/workspace-canvas-plan/01_PLAN.md). Operates on files
inside the session's storage, with the files-workspace draft convention baked
in: editing a file that is NOT already a draft first copies it to
``workspace/drafts/<stem>/<name>`` and edits the copy — originals under
``workspace/uploads/`` are never mutated. After each edit a preview
(``workspace/drafts/<stem>/preview/page-N.png``) is regenerated best-effort via
LibreOffice headless + pdftoppm (both ship in the backend image; see README).
Deliver results to the user with SendUserFile.

Path arguments are relative to the session storage root (e.g.
``workspace/uploads/deck.pptx``) or absolute paths inside it.
"""

from __future__ import annotations

import json
import shutil
import subprocess
import threading
from logging import getLogger
from pathlib import Path
from typing import Any, Dict, List, Optional

from tools.base import BaseTool, ToolError

logger = getLogger(__name__)

# LibreOffice headless is not concurrency-safe (shared profile) — serialize.
_SOFFICE_LOCK = threading.Lock()
_SOFFICE_TIMEOUT = 120
_PREVIEW_DPI = "96"
_MAX_PREVIEW_PAGES = 30


# ── session storage helpers ─────────────────────────────────────────


def _storage_root(session_id: str) -> Path:
    from service.executor import get_agent_session_manager

    manager = get_agent_session_manager()
    agent = manager.get_agent(session_id)
    storage = getattr(agent, "storage_path", None) if agent else None
    if not storage:
        try:
            from service.sessions.store import get_session_store

            rec = get_session_store().get(session_id)
            storage = (rec or {}).get("storage_path")
        except Exception:  # noqa: BLE001
            storage = None
    if not storage:
        raise ToolError("Session storage is not available for this session.")
    return Path(storage).resolve()


def _resolve(root: Path, path: str) -> Path:
    p = Path(path)
    target = (p if p.is_absolute() else root / p).resolve()
    try:
        target.relative_to(root)
    except ValueError:
        raise ToolError(f"Path escapes the session storage: {path}")
    return target


def _rel(root: Path, p: Path) -> str:
    return p.resolve().relative_to(root).as_posix()


def _ensure_draft(root: Path, src: Path) -> Path:
    """Return the working copy for ``src`` (drafts convention).

    Files already under ``workspace/drafts/`` are edited in place; anything
    else is copied to ``workspace/drafts/<stem>/<name>`` first so originals
    (esp. user uploads) stay pristine.
    """
    drafts_root = root / "workspace" / "drafts"
    try:
        src.relative_to(drafts_root)
        return src
    except ValueError:
        pass
    draft_dir = drafts_root / src.stem
    draft_dir.mkdir(parents=True, exist_ok=True)
    draft = draft_dir / src.name
    if not draft.exists():
        shutil.copy2(src, draft)
    return draft


# ── conversion / preview (LibreOffice + poppler) ────────────────────


def _find_soffice() -> Optional[str]:
    """PATH first, then the Debian slim install location (with
    --no-install-recommends the /usr/bin/soffice symlink is absent)."""
    found = shutil.which("soffice")
    if found:
        return found
    fallback = "/usr/lib/libreoffice/program/soffice"
    return fallback if Path(fallback).exists() else None


def _soffice(args: List[str], cwd: Path) -> None:
    binary = _find_soffice()
    if binary is None:
        raise ToolError(
            "LibreOffice (soffice) is not installed in this backend — document "
            "conversion/preview is unavailable. See README (recommended dependency)."
        )
    with _SOFFICE_LOCK:
        proc = subprocess.run(
            [binary, "--headless", "--norestore", *args],
            cwd=str(cwd), capture_output=True, timeout=_SOFFICE_TIMEOUT,
        )
    if proc.returncode != 0:
        raise ToolError(f"LibreOffice conversion failed: {proc.stderr.decode(errors='replace')[:400]}")


def _to_pdf(doc: Path, outdir: Path) -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    _soffice(["--convert-to", "pdf", "--outdir", str(outdir), str(doc)], cwd=outdir)
    pdf = outdir / (doc.stem + ".pdf")
    if not pdf.exists():
        raise ToolError("Conversion produced no PDF output.")
    return pdf


def _pdf_to_pngs(pdf: Path, outdir: Path) -> List[Path]:
    if shutil.which("pdftoppm") is None:
        raise ToolError("pdftoppm (poppler-utils) is not installed — PNG preview unavailable.")
    outdir.mkdir(parents=True, exist_ok=True)
    for old in outdir.glob("page-*.png"):
        old.unlink(missing_ok=True)
    subprocess.run(
        ["pdftoppm", "-png", "-r", _PREVIEW_DPI, "-l", str(_MAX_PREVIEW_PAGES),
         str(pdf), str(outdir / "page")],
        capture_output=True, timeout=_SOFFICE_TIMEOUT, check=True,
    )
    return sorted(outdir.glob("page-*.png"))


def _regen_preview(root: Path, draft: Path) -> Dict[str, Any]:
    """Best-effort preview regeneration for a draft; never raises."""
    try:
        preview_dir = draft.parent / "preview"
        pdf = _to_pdf(draft, preview_dir)
        pages = _pdf_to_pngs(pdf, preview_dir)
        return {"preview_pages": [_rel(root, p) for p in pages]}
    except Exception as exc:  # noqa: BLE001
        return {"preview_error": str(exc)[:200]}


# ── tools ────────────────────────────────────────────────────────────


class DocConvertTool(BaseTool):
    """Convert an office document (pptx/docx/xlsx/…) to pdf, png previews, or text."""

    name = "doc_convert"
    description = (
        "Convert a document in the session storage (pptx/docx/xlsx/odt/…) to another "
        "format. to='pdf' produces a PDF next to a draft copy; to='png' renders "
        "page/slide preview images; to='text' extracts plain text (read the result "
        "with Read). Paths are relative to the session storage, e.g. "
        "'workspace/uploads/deck.pptx'."
    )
    CAPABILITIES_CONCURRENCY_SAFE = False

    def run(self, session_id: str, path: str, to: str = "pdf") -> str:
        """Convert a document.

        Args:
            path: Source file path (relative to session storage or absolute inside it).
            to: Target — 'pdf', 'png' (per-page preview images), or 'text'.
        """
        root = _storage_root(session_id)
        src = _resolve(root, path)
        if not src.is_file():
            raise ToolError(f"File not found: {path}")
        to = (to or "pdf").lower().strip()
        draft = _ensure_draft(root, src)
        outdir = draft.parent

        if to == "pdf":
            pdf = _to_pdf(draft, outdir)
            return json.dumps({"ok": True, "pdf": _rel(root, pdf)}, ensure_ascii=False)
        if to == "png":
            pdf = _to_pdf(draft, outdir / "preview")
            pages = _pdf_to_pngs(pdf, outdir / "preview")
            return json.dumps(
                {"ok": True, "pages": [_rel(root, p) for p in pages]}, ensure_ascii=False
            )
        if to in ("text", "txt"):
            txt = self._extract_text(draft)
            out = outdir / (draft.stem + ".txt")
            out.write_text(txt, encoding="utf-8")
            return json.dumps(
                {"ok": True, "text_file": _rel(root, out), "chars": len(txt),
                 "head": txt[:600]}, ensure_ascii=False
            )
        raise ToolError(f"Unsupported target format: {to} (use pdf / png / text)")

    @staticmethod
    def _extract_text(doc: Path) -> str:
        suffix = doc.suffix.lower()
        if suffix == ".pptx":
            from pptx import Presentation

            parts: List[str] = []
            prs = Presentation(str(doc))
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"--- slide {i} ---")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)
        if suffix == ".docx":
            from docx import Document

            return "\n".join(p.text for p in Document(str(doc)).paragraphs)
        if suffix in (".xlsx", ".xlsm"):
            from openpyxl import load_workbook

            wb = load_workbook(str(doc), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"--- sheet {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join("" if v is None else str(v) for v in row))
            return "\n".join(parts)
        # Fall back to LibreOffice for other formats (odt, doc, ppt, …).
        outdir = doc.parent
        _soffice(["--convert-to", "txt:Text", "--outdir", str(outdir), str(doc)], cwd=outdir)
        out = outdir / (doc.stem + ".txt")
        return out.read_text(encoding="utf-8", errors="replace") if out.exists() else ""


class PptxEditTool(BaseTool):
    """Edit a PowerPoint (pptx): replace text, set slide text, add/delete slides."""

    name = "pptx_edit"
    description = (
        "Edit a .pptx in the session storage. Edits a DRAFT copy under "
        "workspace/drafts/ (originals stay untouched) and regenerates PNG previews. "
        "operations is a JSON list; supported ops: "
        '{"op":"replace_text","find":"…","replace":"…"} (all slides), '
        '{"op":"set_slide_text","slide":1,"placeholder":0,"text":"…"} (1-based slide, '
        "placeholder index), "
        '{"op":"add_slide","title":"…","body":"…"}, '
        '{"op":"delete_slide","slide":2}. '
        "Returns the draft path + preview image paths. Send the finished file to the "
        "user with SendUserFile."
    )

    def run(self, session_id: str, path: str, operations: List[Dict[str, Any]]) -> str:
        """Apply edit operations to a pptx.

        Args:
            path: Source .pptx path (relative to session storage).
            operations: List of operation objects (see tool description).
        """
        from pptx import Presentation
        from pptx.util import Inches

        root = _storage_root(session_id)
        src = _resolve(root, path)
        if not src.is_file():
            raise ToolError(f"File not found: {path}")
        if src.suffix.lower() != ".pptx":
            raise ToolError("pptx_edit only handles .pptx files (use doc_convert for others).")
        draft = _ensure_draft(root, src)
        prs = Presentation(str(draft))
        applied: List[str] = []

        for op in operations or []:
            kind = (op.get("op") or "").strip()
            if kind == "replace_text":
                find, repl = str(op.get("find", "")), str(op.get("replace", ""))
                if not find:
                    raise ToolError("replace_text needs a non-empty 'find'.")
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not getattr(shape, "has_text_frame", False):
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if find in run.text:
                                    run.text = run.text.replace(find, repl)
                                    count += 1
                applied.append(f"replace_text '{find}'→'{repl}' ({count} runs)")
            elif kind == "set_slide_text":
                idx = int(op.get("slide", 1)) - 1
                slides = list(prs.slides)
                if not (0 <= idx < len(slides)):
                    raise ToolError(f"slide {idx + 1} out of range (1..{len(slides)})")
                ph_idx = int(op.get("placeholder", 0))
                phs = [s for s in slides[idx].shapes if getattr(s, "has_text_frame", False)]
                if not (0 <= ph_idx < len(phs)):
                    raise ToolError(f"placeholder {ph_idx} out of range (0..{len(phs) - 1})")
                phs[ph_idx].text_frame.text = str(op.get("text", ""))
                applied.append(f"set_slide_text slide {idx + 1} ph {ph_idx}")
            elif kind == "add_slide":
                layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title is not None:
                    slide.shapes.title.text = str(op.get("title", ""))
                body = str(op.get("body", ""))
                if body:
                    placed = False
                    for shape in slide.placeholders:
                        if shape != slide.shapes.title and getattr(shape, "has_text_frame", False):
                            shape.text_frame.text = body
                            placed = True
                            break
                    if not placed:
                        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                        box.text_frame.text = body
                applied.append("add_slide")
            elif kind == "delete_slide":
                idx = int(op.get("slide", 1)) - 1
                xml_slides = prs.slides._sldIdLst  # noqa: SLF001 — no public delete API
                ids = list(xml_slides)
                if not (0 <= idx < len(ids)):
                    raise ToolError(f"slide {idx + 1} out of range (1..{len(ids)})")
                xml_slides.remove(ids[idx])
                applied.append(f"delete_slide {idx + 1}")
            else:
                raise ToolError(f"Unknown op: {kind!r}")

        prs.save(str(draft))
        result: Dict[str, Any] = {
            "ok": True, "draft": _rel(root, draft), "applied": applied,
            "slides": len(list(prs.slides)),
        }
        result.update(_regen_preview(root, draft))
        return json.dumps(result, ensure_ascii=False)


class XlsxEditTool(BaseTool):
    """Edit an Excel workbook (xlsx): set cells, add sheets."""

    name = "xlsx_edit"
    description = (
        "Edit a .xlsx in the session storage. Edits a DRAFT copy under "
        "workspace/drafts/ (originals stay untouched). operations is a JSON list; "
        'supported ops: {"op":"set_cell","sheet":"Sheet1","cell":"B2","value":123}, '
        '{"op":"add_sheet","name":"Data"}. value type is preserved '
        "(number/string/bool). Returns the draft path."
    )

    def run(self, session_id: str, path: str, operations: List[Dict[str, Any]]) -> str:
        """Apply edit operations to an xlsx.

        Args:
            path: Source .xlsx path (relative to session storage).
            operations: List of operation objects (see tool description).
        """
        from openpyxl import load_workbook

        root = _storage_root(session_id)
        src = _resolve(root, path)
        if not src.is_file():
            raise ToolError(f"File not found: {path}")
        if src.suffix.lower() not in (".xlsx", ".xlsm"):
            raise ToolError("xlsx_edit only handles .xlsx/.xlsm files.")
        draft = _ensure_draft(root, src)
        wb = load_workbook(str(draft))
        applied: List[str] = []

        for op in operations or []:
            kind = (op.get("op") or "").strip()
            if kind == "set_cell":
                sheet = op.get("sheet") or wb.sheetnames[0]
                if sheet not in wb.sheetnames:
                    raise ToolError(f"Sheet not found: {sheet} (have: {', '.join(wb.sheetnames)})")
                cell = str(op.get("cell", "")).strip()
                if not cell:
                    raise ToolError("set_cell needs a 'cell' (e.g. 'B2').")
                wb[sheet][cell] = op.get("value")
                applied.append(f"set_cell {sheet}!{cell}")
            elif kind == "add_sheet":
                name = str(op.get("name", "")).strip() or f"Sheet{len(wb.sheetnames) + 1}"
                wb.create_sheet(title=name)
                applied.append(f"add_sheet {name}")
            else:
                raise ToolError(f"Unknown op: {kind!r}")

        wb.save(str(draft))
        return json.dumps(
            {"ok": True, "draft": _rel(root, draft), "applied": applied,
             "sheets": wb.sheetnames}, ensure_ascii=False
        )


class DocxEditTool(BaseTool):
    """Edit a Word document (docx): replace text, append paragraphs/headings."""

    name = "docx_edit"
    description = (
        "Edit a .docx in the session storage. Edits a DRAFT copy under "
        "workspace/drafts/ (originals stay untouched) and regenerates PNG previews. "
        "operations is a JSON list; supported ops: "
        '{"op":"replace_text","find":"…","replace":"…"}, '
        '{"op":"append_paragraph","text":"…","heading":0} (heading 1-9 for a heading, '
        "0/omitted for body text). Returns the draft path + preview image paths."
    )

    def run(self, session_id: str, path: str, operations: List[Dict[str, Any]]) -> str:
        """Apply edit operations to a docx.

        Args:
            path: Source .docx path (relative to session storage).
            operations: List of operation objects (see tool description).
        """
        from docx import Document

        root = _storage_root(session_id)
        src = _resolve(root, path)
        if not src.is_file():
            raise ToolError(f"File not found: {path}")
        if src.suffix.lower() != ".docx":
            raise ToolError("docx_edit only handles .docx files (use doc_convert for others).")
        draft = _ensure_draft(root, src)
        doc = Document(str(draft))
        applied: List[str] = []

        for op in operations or []:
            kind = (op.get("op") or "").strip()
            if kind == "replace_text":
                find, repl = str(op.get("find", "")), str(op.get("replace", ""))
                if not find:
                    raise ToolError("replace_text needs a non-empty 'find'.")
                count = 0
                for para in doc.paragraphs:
                    for run in para.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, repl)
                            count += 1
                applied.append(f"replace_text '{find}'→'{repl}' ({count} runs)")
            elif kind == "append_paragraph":
                text = str(op.get("text", ""))
                level = int(op.get("heading", 0) or 0)
                if level > 0:
                    doc.add_heading(text, level=min(level, 9))
                else:
                    doc.add_paragraph(text)
                applied.append("append_paragraph")
            else:
                raise ToolError(f"Unknown op: {kind!r}")

        doc.save(str(draft))
        result: Dict[str, Any] = {"ok": True, "draft": _rel(root, draft), "applied": applied}
        result.update(_regen_preview(root, draft))
        return json.dumps(result, ensure_ascii=False)


TOOLS = [DocConvertTool(), PptxEditTool(), XlsxEditTool(), DocxEditTool()]
